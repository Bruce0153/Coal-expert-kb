from __future__ import annotations

from typing import List

from langchain_core.documents import Document

from coal_kb.loaders.base import BaseLoader, detect_language, normalize_text
from coal_kb.loaders.registry import register_loader


class PptxLoader(BaseLoader):
    def __init__(self) -> None:
        try:
            import pptx  # noqa: F401
        except Exception as exc:  # pragma: no cover - optional dependency
            raise RuntimeError("python-pptx is required for PptxLoader") from exc

    def load(self, path: str) -> List[Document]:
        from pptx import Presentation

        prs = Presentation(path)
        docs: List[Document] = []
        for idx, slide in enumerate(prs.slides):
            parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    parts.append(shape.text)
            content = normalize_text(" ".join(parts))
            if not content:
                continue
            lang = detect_language(content)
            docs.append(
                Document(
                    page_content=content,
                    metadata={
                        "source_file": path,
                        "section": "slide",
                        "page": idx,
                        "doc_type": "pptx",
                        "language": lang,
                        "parser": "pptx",
                    },
                )
            )
        return docs


register_loader("pptx", PptxLoader)
